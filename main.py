import PyPDF2
from PyPDF2 import PdfReader
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from docx2pdf import convert
from openai import OpenAI
from pydub import AudioSegment
import os

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

system_prompt = ""

def pdf_to_text(pdf_file_path):
    text = ""
    with open(pdf_file_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
            
    return text


def pptx_to_pdf(pptx_file, pdf_file):
    presentation = Presentation(pptx_file)
    
    pdf_canvas = canvas.Canvas(pdf_file, pagesize=letter)
    
    for slide in presentation.slides:
        pdf_canvas.showPage()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pdf_canvas.drawString(shape.left, shape.top - run.font.size, run.text)
    return pdf_canvas


def pptx_to_text(pptx_file):
    base_filename = os.path.splitext(os.path.basename(pptx_file))[0]
    output_filename = f"{base_filename}.pdf"
    output_canvas = pptx_to_pdf(pptx_file, output_filename)
    text = pdf_to_text(output_canvas)
    return text

def text_to_enhanced_text(text):
    text = text
    response = client.chat.completions.create(
        model="gpt-4-1106-preview",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": text}
        ]
    )
    return response.choices[0].message.content



def text_to_speech(text_file):
    try:
        base_filename = os.path.splitext(os.path.basename(text_file))[0]
        with open(text_file, 'r', encoding='utf-8') as file:
            file_contents = file.read()
        
        chunk_size = 4000
        text_chunks = [file_contents[i:i+chunk_size] for i in range(0, len(file_contents), chunk_size)]
        mp3_files = []
        
        for i, chunk in enumerate(text_chunks):
            response = client.audio.speech.create(
                model="tts-1",
                voice="onyx",
                input=chunk,
            )
            mp3_filename = f"chunk_{i}.mp3"
            response.stream_to_file(mp3_filename)
            print(f"Generated {mp3_filename}")
            mp3_files.append(mp3_filename)

        # Combine the MP3 files into one
        combined = AudioSegment.silent()
        for mp3_file in mp3_files:
            combined += AudioSegment.from_mp3(mp3_file)
        combined.export(f"{base_filename}.mp3", format="mp3")
        print("Combined MP3 files into 'output.mp3'")
        for mp3_file in mp3_files:
            os.remove(mp3_file)
        print("Deleted temporary MP3 files")
    except FileNotFoundError:
        print(f"File not found: {text_file}")
    except Exception as e:
        print(f"An error occurred: {str(e)}")



if __name__ == "__main__":
    
    system_prompt = input("Enter the system prompt: ")
    for filename in os.listdir():
        if filename.endswith(".pdf"):
            extracted_text = pdf_to_text(filename)
        # elif filename.endswith(".pptx"):
        #     extracted_text = pptx_to_text(filename)
        else:
            print("Invalid file format")
            continue
        
        enhanced_text = text_to_enhanced_text(extracted_text)
        base_filename = os.path.splitext(os.path.basename(filename))[0]
        output_filename = f"{base_filename}.txt"
        
        with open(output_filename, "w", encoding="utf-8") as output_file:
            output_file.write(enhanced_text)

        print(f"Text extracted and saved to {output_filename}")
        text_to_speech(output_filename)
    